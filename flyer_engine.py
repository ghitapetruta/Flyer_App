"""
flyer_engine.py — motorul de generare HTML și PPTX pentru flyere.

Acest modul conține funcțiile esențiale extrase din generator_flyer_v10.py,
adaptate pentru a primi date direct (în loc de a citi din Excel) și a
returna conținutul ca bytes (în loc de a scrie pe disk).

Folosit de aplicația Streamlit (app.py).
"""
from __future__ import annotations
import base64
import io
from pathlib import Path
from typing import Optional


# ══════════════════════════════════════════════════════════════
# DEFAULT STYLE
# ══════════════════════════════════════════════════════════════
DEFAULT_STYLE = {
    "Culoare principală (roșu)":          "A11319",
    "Culoare secundară (gri)":            "565757",
    "Culoare fundal pagină":              "FFFFFF",
    "Culoare fundal casetă produs":       "F9F9F7",
    "Culoare fundal pill caracteristici": "EEEDE9",
    "Culoare bordură":                    "C8C7C3",
    "Culoare footer (jos pagină)":        "1C1C1C",
    "Font categorie (sus)":               "Montserrat",
    "Dimensiune categorie (sus)":         9,
    "Culoare categorie (sus)":            "BBBBBB",
    "Font titlu produs (mare)":           "Montserrat",
    "Dimensiune titlu produs (mare)":     28,
    "Culoare titlu produs (mare)":        "FFFFFF",
    "Font text promoțional mic":          "Montserrat",
    "Dimensiune text promoțional mic":    9,
    "Culoare text promoțional mic":       "FFCCCC",
    "Font text promoțional mare":         "Montserrat",
    "Dimensiune text promoțional mare":   17,
    "Culoare text promoțional mare":      "FFFFFF",
    "Font denumire produs":               "Montserrat",
    "Dimensiune denumire produs":         18,
    "Culoare denumire produs":            "1C1C1C",
    "Font caracteristici":                "Arial",
    "Dimensiune caracteristici":          13,
    "Culoare caracteristici (etichetă)":  "333333",
    "Culoare caracteristici (valoare)":   "1C1C1C",
    "Font cod produs":                    "Courier New",
    "Dimensiune cod produs":              13,
    "Culoare cod produs":                 "555555",
    "Font preț":                          "Montserrat",
    "Dimensiune preț":                    32,
    "Culoare preț":                       "FFFFFF",
    "Dimensiune »LEI« preț":              17,
    "Culoare »LEI« preț":                 "FFE0E0",
    "Dimensiune UM (/ buc)":              14,
    "Culoare UM (/ buc)":                 "FFCCCC",
    "Dimensiune preț vechi (tăiat)":      16,
    "Culoare preț vechi (tăiat)":         "FFB0B0",
    "Dimensiune »Preț cu TVA«":           13,
    "Culoare »Preț cu TVA«":              "FFD0D0",
    "Font observații":                    "Arial",
    "Dimensiune observații":              12,
    "Culoare titlu observații":           "1C1C1C",
    "Culoare text observații":            "333333",
    "Culoare nume firmă (footer)":        "1C1C1C",
    "Diametru bulină reducere (px)":      56,
    "Dimensiune cifră reducere":          24,
    "Font badge":                         "Arial",
    "Culoare fundal badge":               "1C1C1C",
    "Culoare text badge":                 "FFFFFF",
    "Dimensiune font badge (px)":         11,
    "Padding badge (px)":                 4,
    "Dimensiune imagine badge (px)":      18,
}

DEFAULT_ANTET = {
    "Nume companie": "SC METALCOM SRL",
    "Categorie (sus)": "ACCESORII MOBILIER · SYSKOR",
    "Titlu produs (mare)": "SERTARE VISION TOP",
    "Text promoțional mic": "OFERTĂ SPECIALĂ",
    "Text promoțional mare": "Prețuri cu TVA inclus",
    "Telefon": "0745 751 235",
    "Email": "comenzi@metalcom.ro",
    "Website 1": "www.metalcom.ro",
    "Website 2": "www.metalcom-shop.ro",
    "Footer legal": "SC Metalcom SRL · CUI RO6778720 · J16/3822/1994 · Str. Calea București nr.191, Craiova, Dolj",
}

DEFAULT_OBSERVATII = [
    "Prețurile sunt exprimate în LEI și conțin TVA.",
    "Reducerile sunt valabile în limita stocului disponibil.",
    "Ofertă valabilă până la epuizarea stocului.",
]


# ══════════════════════════════════════════════════════════════
# IMAGE PROCESSING
# ══════════════════════════════════════════════════════════════
def img_bytes_to_b64(img_bytes: bytes, max_width: int = 800) -> str:
    """Convertește bytes → base64 data URL, redimensionând la nevoie."""
    if not img_bytes:
        return ""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(img_bytes))
        if img.mode != "RGB":
            img = img.convert("RGB")
        if img.width > max_width:
            ratio = max_width / img.width
            img = img.resize((max_width, int(img.height * ratio)), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode()
    except Exception:
        # fallback - returnăm bytes brute fără redimensionare
        return "data:image/jpeg;base64," + base64.b64encode(img_bytes).decode()


def hex_to_rgb(hex_color: str):
    s = (hex_color or "").lstrip("#").strip().upper()
    if len(s) != 6:
        return (0, 0, 0)
    try:
        return tuple(int(s[i:i+2], 16) for i in (0, 2, 4))
    except ValueError:
        return (0, 0, 0)


def normalize_color(val):
    if val is None:
        return ""
    return str(val).strip().lstrip("#").upper()


# ══════════════════════════════════════════════════════════════
# LAYOUT HELPERS
# ══════════════════════════════════════════════════════════════
def _products_per_page(per_page) -> int:
    if per_page == "3a":
        return 3
    return int(per_page)


def _grid_layout_css(per_page) -> str:
    layouts = {
        1: "grid-template-columns: 1fr;\n  grid-template-rows: minmax(0, 1fr);",
        2: "grid-template-columns: 1fr;\n  grid-template-rows: repeat(2, minmax(0, 1fr));",
        3: "grid-template-columns: 1fr;\n  grid-template-rows: repeat(3, minmax(0, 1fr));",
        "3a": "grid-template-columns: 1fr 1fr;\n  grid-template-rows: repeat(3, minmax(0, 1fr));",
        4: "grid-template-columns: 1fr 1fr;\n  grid-template-rows: minmax(0, 1fr) minmax(0, 1fr);",
        5: "grid-template-columns: 1fr 1fr;\n  grid-template-rows: minmax(0, 1fr) minmax(0, 1fr) minmax(0, 1fr);",
        6: "grid-template-columns: 1fr 1fr;\n  grid-template-rows: repeat(3, minmax(0, 1fr));",
    }
    return layouts.get(per_page, layouts[6])


def _special_card_css(per_page) -> str:
    if per_page == "3a":
        return """
.grid > .card:nth-child(1) { grid-column: span 2; grid-row: span 2; }
.grid > .card:nth-child(1) .card-name { font-size: 26px; }
.grid > .card:nth-child(1) .card-img-wrap img { max-height: 280px; }
.grid > .card:nth-child(1) .card-bottom { min-height: 88px; height: 88px; }
.grid > .card:nth-child(1) .card-pret,
.grid > .card:nth-child(1) .pret-nou { font-size: 48px; }
.grid > .card:nth-child(1) .lei { font-size: 24px; }
.grid > .card:nth-child(1) .pret-vechi { font-size: 20px; }
.grid > .card:nth-child(1) .um { font-size: 20px; }
.grid > .card:nth-child(1) .card-tva { font-size: 16px; }
.grid > .card:nth-child(1) .pill { font-size: 16px; padding: 5px 12px; }
.grid > .card:nth-child(1) .card-cod { font-size: 16px; }
.grid > .card:nth-child(1) .product-badge { font-size: 14px; padding: 5px 13px; }
.grid > .card:nth-child(1) .product-badge .badge-icon { height: 24px; }
.grid > .card:nth-child(1) .discount-badge { width: 80px; height: 80px; }
.grid > .card:nth-child(1) .discount-badge .dscV { font-size: 32px; }
.grid > .card:nth-child(1) .discount-badge .dscP { font-size: 18px; }
.grid > .card:nth-child(2) .card-name,
.grid > .card:nth-child(3) .card-name { font-size: 16px; }
.grid > .card:nth-child(2) .card-bottom,
.grid > .card:nth-child(3) .card-bottom { min-height: 56px; height: 56px; }
.grid > .card:nth-child(2) .card-pret,
.grid > .card:nth-child(2) .pret-nou,
.grid > .card:nth-child(3) .card-pret,
.grid > .card:nth-child(3) .pret-nou { font-size: 28px; }
.grid > .card:nth-child(2) .lei,
.grid > .card:nth-child(3) .lei { font-size: 14px; }
"""
    elif per_page == 5:
        return ".grid > .card:nth-child(5) { grid-column: span 2; }\n"
    elif per_page == 1:
        return """
.grid > .card .card-img-wrap { width: 50%; padding: 24px; }
.grid > .card .card-img-wrap img { max-height: none; }
.grid > .card .card-info { padding: 30px 32px 24px; gap: 16px; }
.grid > .card .card-name { font-size: 32px; }
.grid > .card .card-cod { font-size: 17px; }
.grid > .card .pill { font-size: 17px; padding: 6px 14px; }
.grid > .card .product-badge { font-size: 16px; padding: 6px 14px; }
.grid > .card .product-badge .badge-icon { height: 26px; }
.grid > .card .card-bottom { padding: 18px 28px; min-height: 100px; height: 100px; }
.grid > .card .card-pret, .grid > .card .pret-nou { font-size: 56px; }
.grid > .card .lei { font-size: 28px; }
.grid > .card .pret-vechi { font-size: 24px; }
.grid > .card .um { font-size: 22px; }
.grid > .card .card-tva { font-size: 18px; }
.grid > .card .discount-badge { width: 90px; height: 90px; top: 12px; left: 12px; }
.grid > .card .discount-badge .dscV { font-size: 36px; }
.grid > .card .discount-badge .dscP { font-size: 20px; }
"""
    elif per_page == 2:
        return """
.grid > .card .card-name { font-size: 24px; }
.grid > .card .card-img-wrap img { max-height: 220px; }
.grid > .card .card-bottom { min-height: 80px; height: 80px; }
.grid > .card .card-pret, .grid > .card .pret-nou { font-size: 42px; }
.grid > .card .lei { font-size: 22px; }
.grid > .card .pret-vechi { font-size: 18px; }
.grid > .card .um { font-size: 18px; }
.grid > .card .card-tva { font-size: 15px; }
.grid > .card .pill { font-size: 15px; padding: 4px 11px; }
.grid > .card .card-cod { font-size: 15px; }
.grid > .card .product-badge { font-size: 13px; padding: 4px 12px; }
.grid > .card .product-badge .badge-icon { height: 22px; }
.grid > .card .discount-badge { width: 70px; height: 70px; }
.grid > .card .discount-badge .dscV { font-size: 28px; }
"""
    elif per_page == 3:
        return """
.grid > .card .card-name { font-size: 20px; }
.grid > .card .card-img-wrap img { max-height: 180px; }
.grid > .card .card-bottom { min-height: 72px; height: 72px; }
.grid > .card .card-pret, .grid > .card .pret-nou { font-size: 38px; }
.grid > .card .lei { font-size: 19px; }
.grid > .card .pret-vechi { font-size: 17px; }
.grid > .card .um { font-size: 16px; }
.grid > .card .pill { font-size: 14px; }
.grid > .card .card-cod { font-size: 14px; }
"""
    elif per_page == 4:
        return """
.grid > .card .card-name { font-size: 22px; }
.grid > .card .card-img-wrap img { max-height: 200px; }
.grid > .card .card-bottom { min-height: 76px; height: 76px; }
.grid > .card .card-pret, .grid > .card .pret-nou { font-size: 40px; }
.grid > .card .lei { font-size: 20px; }
.grid > .card .pret-vechi { font-size: 17px; }
.grid > .card .um { font-size: 16px; }
.grid > .card .pill { font-size: 14px; }
.grid > .card .card-cod { font-size: 14px; }
"""
    return ""


# ══════════════════════════════════════════════════════════════
# CSS GENERATION
# ══════════════════════════════════════════════════════════════
def build_css(s: dict, per_page=6) -> str:
    def color(key): return "#" + (s.get(key) or DEFAULT_STYLE[key]).lstrip("#")
    def fs(key):    return f"{s.get(key, DEFAULT_STYLE[key])}px"
    def font(key):  return f"'{s.get(key, DEFAULT_STYLE[key])}', sans-serif"

    fonts_used = {str(s.get(k, DEFAULT_STYLE[k])).strip() for k in s if k.startswith("Font ")}
    google_fonts = {
        "Montserrat", "Roboto", "Open Sans", "Lato", "Poppins",
        "Oswald", "Raleway", "Source Sans Pro", "Ubuntu", "Nunito",
        "Inter", "Playfair Display", "Merriweather", "PT Sans",
        "Barlow", "Barlow Condensed", "Work Sans", "Rubik", "DM Sans",
    }
    web_fonts_to_load = [f for f in fonts_used if f in google_fonts]
    google_import = ""
    if web_fonts_to_load:
        font_query = "&family=".join(
            f.replace(" ", "+") + ":wght@400;500;600;700;800;900"
            for f in web_fonts_to_load
        )
        google_import = f"@import url('https://fonts.googleapis.com/css2?family={font_query}&display=swap');\n"

    discount_size = int(s.get("Dimensiune cifră reducere", DEFAULT_STYLE["Dimensiune cifră reducere"]))
    badge_pad = int(s.get("Padding badge (px)", DEFAULT_STYLE["Padding badge (px)"]))

    return f"""
{google_import}
*, *::before, *::after {{
  box-sizing: border-box; margin: 0; padding: 0;
  -webkit-print-color-adjust: exact !important;
  print-color-adjust: exact !important;
  color-adjust: exact !important;
}}

body {{
  font-family: 'Arial', sans-serif;
  background: #888;
  padding: 24px 16px;
  display: flex; flex-direction: column; align-items: center; gap: 24px;
}}

.page {{
  width: 210mm;
  height: 297mm;
  background: {color('Culoare fundal pagină')} !important;
  box-shadow: 0 8px 48px rgba(0,0,0,.28);
  display: flex; flex-direction: column;
  overflow: hidden; page-break-after: always;
}}

.header {{ border-bottom: 3px solid {color('Culoare principală (roșu)')}; line-height: 0; flex-shrink: 0; }}
.header img {{ width: 100%; display: block; }}

.hero {{
  background-color: {color('Culoare secundară (gri)')} !important;
  padding: 12px 22px;
  display: flex; align-items: center; justify-content: space-between;
  position: relative; overflow: hidden;
  flex-shrink: 0;
}}
.hero::after {{
  content: ''; position: absolute; top: 0; right: 0;
  width: 155px; height: 100%;
  background-color: {color('Culoare principală (roșu)')} !important;
  clip-path: polygon(28% 0, 100% 0, 100% 100%, 0% 100%);
}}
.hero-left {{ position: relative; z-index: 1; }}
.hero-cat {{
  font-family: {font('Font categorie (sus)')};
  font-size: {fs('Dimensiune categorie (sus)')};
  letter-spacing: 3px; text-transform: uppercase;
  color: {color('Culoare categorie (sus)')} !important;
}}
.hero-title {{
  font-family: {font('Font titlu produs (mare)')};
  font-size: {fs('Dimensiune titlu produs (mare)')};
  font-weight: 900;
  color: {color('Culoare titlu produs (mare)')} !important;
  text-transform: uppercase; line-height: 1; letter-spacing: .5px;
}}
.hero-right {{ position: relative; z-index: 1; text-align: right; }}
.hl1 {{
  font-family: {font('Font text promoțional mic')};
  font-size: {fs('Dimensiune text promoțional mic')};
  letter-spacing: 2px; text-transform: uppercase;
  color: {color('Culoare text promoțional mic')} !important;
  display: block;
}}
.hl2 {{
  font-family: {font('Font text promoțional mare')};
  font-size: {fs('Dimensiune text promoțional mare')};
  font-weight: 700;
  color: {color('Culoare text promoțional mare')} !important;
  display: block;
}}

.grid {{
  display: grid;
  {_grid_layout_css(per_page)}
  gap: 9px; padding: 9px 11px;
  flex: 1 1 auto;
  min-height: 0;
  overflow: hidden;
}}

{_special_card_css(per_page)}

.card {{
  display: flex; flex-direction: column;
  border: 1.5px solid {color('Culoare bordură')};
  background-color: {color('Culoare fundal casetă produs')} !important;
  overflow: hidden;
}}
.card-top {{ display: flex; flex-direction: row; flex: 1; }}
.card-img-wrap {{
  width: 46%; flex-shrink: 0;
  background-color: #fff !important;
  border-right: 1px solid {color('Culoare bordură')};
  display: flex; align-items: center; justify-content: center;
  padding: 10px 8px;
  position: relative;
}}
.card-img-wrap img {{
  width: 100%; height: 100%;
  object-fit: contain; display: block; max-height: 145px;
}}

.discount-badge {{
  position: absolute; top: 6px; left: 6px;
  width: {s.get("Diametru bulină reducere (px)", DEFAULT_STYLE["Diametru bulină reducere (px)"])}px;
  height: {s.get("Diametru bulină reducere (px)", DEFAULT_STYLE["Diametru bulină reducere (px)"])}px;
  background-color: {color('Culoare principală (roșu)')} !important;
  color: #fff !important;
  border-radius: 50%;
  display: flex; align-items: center; justify-content: center;
  line-height: 1;
  box-shadow: 0 2px 6px rgba(0,0,0,.25);
  transform: rotate(-8deg);
  border: 2px solid #fff;
}}
.discount-badge .dscV {{
  font-size: {fs('Dimensiune cifră reducere')}; font-weight: 900;
  color: #fff !important;
  display: flex; align-items: baseline;
}}
.discount-badge .dscP {{
  font-size: {int(discount_size * 0.6)}px;
  font-weight: 700;
  color: rgba(255,255,255,.9) !important;
  margin-left: 1px;
}}
.discount-badge .dscV-custom {{
  font-size: {int(discount_size * 0.85)}px;
  font-weight: 900;
  color: #fff !important;
  text-align: center;
  line-height: 1;
  letter-spacing: 0.5px;
  padding: 0 4px;
}}
.discount-badge.text-only .dscV-custom {{ text-transform: uppercase; letter-spacing: 1px; }}
.discount-badge.text-only.text-short .dscV-custom {{ font-size: {int(discount_size * 0.75)}px; }}
.discount-badge.text-only.text-medium .dscV-custom {{
  font-size: {int(discount_size * 0.55)}px; letter-spacing: 0.5px;
}}
.discount-badge.text-only.text-long .dscV-custom {{
  font-size: {int(discount_size * 0.42)}px; letter-spacing: 0;
}}

.card-info {{
  flex: 1; padding: 12px 14px 10px;
  display: flex; flex-direction: column; gap: 8px;
}}
.card-name {{
  font-family: {font('Font denumire produs')};
  font-size: {fs('Dimensiune denumire produs')};
  font-weight: 900;
  color: {color('Culoare denumire produs')} !important;
  text-transform: uppercase; letter-spacing: .4px; line-height: 1.15;
}}
.card-specs {{ display: flex; flex-direction: column; gap: 5px; }}
.pill {{
  font-family: {font('Font caracteristici')};
  background-color: {color('Culoare fundal pill caracteristici')} !important;
  border: 1px solid {color('Culoare bordură')};
  font-size: {fs('Dimensiune caracteristici')};
  color: {color('Culoare caracteristici (etichetă)')} !important;
  padding: 3px 9px; display: inline-block; width: fit-content;
}}
.pill b {{
  color: {color('Culoare caracteristici (valoare)')} !important;
  font-weight: 700;
}}
.card-cod {{
  font-family: {font('Font cod produs')};
  font-size: {fs('Dimensiune cod produs')};
  color: {color('Culoare cod produs')} !important;
  letter-spacing: .4px;
}}

.product-badge {{
  display: inline-flex;
  align-items: center;
  gap: 6px;
  font-family: {font('Font badge')};
  background-color: {color('Culoare fundal badge')} !important;
  color: {color('Culoare text badge')} !important;
  font-size: {s.get("Dimensiune font badge (px)", DEFAULT_STYLE["Dimensiune font badge (px)"])}px;
  font-weight: 700;
  letter-spacing: 1.5px;
  text-transform: uppercase;
  padding: {badge_pad}px {badge_pad * 2 + 4}px;
  width: fit-content;
  margin-top: auto;
  align-self: flex-start;
  line-height: 1;
}}
.product-badge .badge-icon {{
  height: {s.get("Dimensiune imagine badge (px)", DEFAULT_STYLE["Dimensiune imagine badge (px)"])}px;
  width: auto;
  display: block;
  object-fit: contain;
}}

.card-bottom {{
  background-color: {color('Culoare principală (roșu)')} !important;
  display: flex; align-items: center; justify-content: space-between;
  padding: 8px 14px;
  border-top: 1px solid rgba(0,0,0,.1);
  min-height: 64px;
  height: 64px;
  box-sizing: border-box;
}}

.card-pret {{
  font-family: {font('Font preț')};
  font-size: {fs('Dimensiune preț')}; font-weight: 900;
  color: {color('Culoare preț')} !important; line-height: 1;
  display: flex; align-items: baseline; gap: 5px;
}}
.lei {{
  font-size: {fs('Dimensiune »LEI« preț')}; font-weight: 600;
  color: {color('Culoare »LEI« preț')} !important;
}}
.um {{
  font-size: {fs('Dimensiune UM (/ buc)')}; font-weight: 500;
  color: {color('Culoare UM (/ buc)')} !important;
  margin-left: 1px;
}}

.card-pret-group {{
  font-family: {font('Font preț')};
  display: flex; flex-direction: column;
  line-height: 1;
  justify-content: center;
}}
.pret-vechi {{
  font-size: {fs('Dimensiune preț vechi (tăiat)')}; font-weight: 500;
  color: {color('Culoare preț vechi (tăiat)')} !important;
  text-decoration: line-through;
  margin-bottom: 3px;
}}
.pret-nou {{
  font-size: {fs('Dimensiune preț')}; font-weight: 900;
  color: {color('Culoare preț')} !important;
  display: flex; align-items: baseline; gap: 4px;
}}
.pret-nou .lei {{
  font-size: {fs('Dimensiune »LEI« preț')}; font-weight: 600;
  color: {color('Culoare »LEI« preț')} !important;
}}
.pret-nou .um {{
  font-size: {fs('Dimensiune UM (/ buc)')}; font-weight: 500;
  color: {color('Culoare UM (/ buc)')} !important;
}}

.card-tva {{
  font-size: {fs('Dimensiune »Preț cu TVA«')};
  font-weight: 600; letter-spacing: .5px;
  text-transform: uppercase;
  color: {color('Culoare »Preț cu TVA«')} !important;
  text-align: right; line-height: 1.3;
}}

.obs {{
  background-color: {color('Culoare fundal pill caracteristici')} !important;
  border-top: 2px solid {color('Culoare bordură')};
  padding: 6px 16px 5px;
  display: flex; align-items: flex-start; justify-content: space-between; gap: 16px;
  flex-shrink: 0;
}}
.obs-title {{
  font-family: {font('Font observații')};
  font-size: {fs('Dimensiune observații')}; font-weight: 700; text-transform: uppercase; letter-spacing: 1.5px;
  color: {color('Culoare titlu observații')} !important; margin-bottom: 2px;
  display: flex; align-items: center; gap: 8px;
}}
.obs-title::before {{
  content: ''; display: inline-block; width: 14px; height: 2px;
  background-color: {color('Culoare principală (roșu)')} !important;
}}
.obs-item {{
  font-family: {font('Font observații')};
  font-size: {fs('Dimensiune observații')};
  color: {color('Culoare text observații')} !important;
  line-height: 1.35; margin-top: 1px;
}}
.obs-right {{
  font-family: {font('Font observații')};
  text-align: right;
  font-size: {fs('Dimensiune observații')};
  color: {color('Culoare text observații')} !important;
  line-height: 1.35;
  border-left: 1px solid {color('Culoare bordură')};
  padding-left: 14px; flex-shrink: 0;
}}
.obs-brand {{
  font-family: {font('Font observații')};
  font-size: {int(s.get("Dimensiune observații", DEFAULT_STYLE["Dimensiune observații"])) + 2}px;
  font-weight: 700; letter-spacing: 1px;
  color: {color('Culoare nume firmă (footer)')} !important; margin-bottom: 2px;
}}

.foot {{
  background-color: {color('Culoare footer (jos pagină)')} !important;
  padding: 5px 16px;
  display: flex; align-items: center; justify-content: space-between;
  flex-shrink: 0;
}}
.foot span {{ font-size: 8px; color: #777 !important; }}

/* Buton flotant pentru Print/Save as PDF */
.print-btn {{
  position: fixed;
  bottom: 24px;
  right: 24px;
  z-index: 9999;
  background: {color('Culoare principală (roșu)')} !important;
  color: #fff !important;
  border: none;
  padding: 14px 22px;
  font-family: 'Arial', sans-serif;
  font-size: 14px;
  font-weight: 700;
  letter-spacing: 0.5px;
  text-transform: uppercase;
  border-radius: 50px;
  box-shadow: 0 4px 16px rgba(0,0,0,.35);
  cursor: pointer;
  transition: transform .15s ease, box-shadow .15s ease;
  display: flex; align-items: center; gap: 8px;
}}
.print-btn:hover {{
  transform: translateY(-2px);
  box-shadow: 0 6px 22px rgba(0,0,0,.45);
}}
.print-btn:active {{ transform: translateY(0); }}
.print-btn::before {{
  content: '⬇';
  font-size: 18px;
  line-height: 1;
}}
.print-hint {{
  position: fixed;
  bottom: 78px;
  right: 24px;
  z-index: 9998;
  background: rgba(28,28,28,.92) !important;
  color: #fff !important;
  padding: 10px 14px;
  font-family: 'Arial', sans-serif;
  font-size: 12px;
  border-radius: 6px;
  max-width: 260px;
  line-height: 1.4;
  box-shadow: 0 4px 12px rgba(0,0,0,.3);
  display: none;
}}
.print-hint.show {{ display: block; }}

@media print {{
  @page {{ margin: 0; size: A4; }}
  body {{ background: none !important; padding: 0; gap: 0; }}
  .page {{ box-shadow: none; width: 210mm; page-break-after: always; }}
  .print-btn, .print-hint {{ display: none !important; }}
}}
"""


# ══════════════════════════════════════════════════════════════
# FORMAT PRICE
# ══════════════════════════════════════════════════════════════
def format_price(val) -> str:
    if val is None or str(val).strip() == "":
        return ""
    try:
        num = float(val)
        if num == int(num):
            return f"{int(num):,}".replace(",", ".")
        return f"{num:,.2f}".replace(",", ".")
    except (ValueError, TypeError):
        return str(val)


# ══════════════════════════════════════════════════════════════
# BUILD CARD HTML
# ══════════════════════════════════════════════════════════════
def build_card(p: dict, img_cache: dict) -> str:
    img_src = img_cache.get(p.get("imagine", ""), "")

    pills_html = ""
    for a in p.get("attrs", []):
        if a.get("label") and a.get("value"):
            pills_html += f'<span class="pill">{a["label"]}: <b>{a["value"]}</b></span>\n        '

    # Bulină
    badge_html = ""
    bulina_text = (p.get("bulina_text") or "").strip()
    if bulina_text:
        looks_like_percent = "%" in bulina_text
        text_len = len(bulina_text)
        size_class = ""
        if not looks_like_percent:
            if text_len <= 3:    size_class = " text-short"
            elif text_len <= 6:  size_class = " text-medium"
            else:                size_class = " text-long"
        css_extra = " text-only" + size_class if not looks_like_percent else ""
        badge_html = f'''<div class="discount-badge{css_extra}">
          <div class="dscV-custom">{bulina_text}</div>
        </div>'''
    else:
        raw_red = p.get("reducere")
        if raw_red is not None and str(raw_red).strip() != "":
            try:
                s = str(raw_red).strip().replace("%", "").replace(",", ".")
                if s:
                    num = float(s)
                    if 0 < num < 1: num *= 100
                    red_val = int(round(num))
                    if red_val > 0:
                        badge_html = f'''<div class="discount-badge">
          <div class="dscV">-{red_val}<span class="dscP">%</span></div>
        </div>'''
            except (ValueError, TypeError):
                pass

    # Preț
    pret_int = format_price(p.get("pret"))
    pret_redus = format_price(p.get("pret_redus"))
    has_pret_redus = (
        p.get("pret_redus") is not None
        and str(p.get("pret_redus")).strip() != ""
        and pret_redus
    )
    moneda = (p.get("moneda") or "").strip()
    um = (p.get("um") or "").strip()
    um_suffix = f' <span class="um">/ {um}</span>' if um else ""
    money_suffix = f' <span class="lei">{moneda}</span>' if moneda else ""

    if has_pret_redus:
        old_money = f' {moneda}' if moneda else ''
        price_html = f'''<div class="card-pret-group">
      <div class="pret-vechi">{pret_int}{old_money}</div>
      <div class="pret-nou">{pret_redus}{money_suffix}{um_suffix}</div>
    </div>'''
    else:
        price_html = f'''<div class="card-pret">{pret_int}{money_suffix}{um_suffix}</div>'''

    # Product badge (text + imagine opționale)
    badge_text = (p.get("badge_text") or "").strip()
    badge_image = (p.get("badge_image") or "").strip()
    product_badge_html = ""
    parts = []
    if badge_image:
        b64 = img_cache.get(badge_image, "")
        if b64:
            parts.append(f'<img class="badge-icon" src="{b64}" alt="">')
    if badge_text:
        parts.append(f'<span class="badge-text">{badge_text}</span>')
    if parts:
        product_badge_html = '<div class="product-badge">' + "".join(parts) + '</div>'

    # TVA text
    tva_text = (p.get("tva_text") or "").strip()
    tva_html = f'<div class="card-tva">{tva_text}</div>' if tva_text else ""

    return f'''<div class="card">
  <div class="card-top">
    <div class="card-img-wrap">
      <img src="{img_src}" alt="{p.get("cod","")}">
      {badge_html}
    </div>
    <div class="card-info">
      <div class="card-name">{p.get("denumire","")}</div>
      <div class="card-cod">COD: {p.get("cod","")}</div>
      <div class="card-specs">
        {pills_html}
      </div>
      {product_badge_html}
    </div>
  </div>
  <div class="card-bottom">
    {price_html}
    {tva_html}
  </div>
</div>'''


# ══════════════════════════════════════════════════════════════
# BUILD PAGE
# ══════════════════════════════════════════════════════════════
def build_page(products, page_num, total_pages, antet, observatii, antet_b64, img_cache):
    cards_html = "\n".join(build_card(p, img_cache) for p in products)
    obs_html = "\n".join(f'<div class="obs-item">▪ {o}</div>' for o in observatii)

    nume    = antet.get("Nume companie", "")
    cat     = antet.get("Categorie (sus)", "")
    titlu   = antet.get("Titlu produs (mare)", "")
    promo_s = antet.get("Text promoțional mic", "")
    promo_m = antet.get("Text promoțional mare", "")
    tel     = antet.get("Telefon", "")
    email   = antet.get("Email", "")
    web1    = antet.get("Website 1", "")
    web2    = antet.get("Website 2", "")
    footer  = antet.get("Footer legal", "")

    header_html = ""
    if antet_b64:
        header_html = f'<div class="header"><img src="{antet_b64}" alt="{nume}"></div>'

    return f'''<div class="page">
  {header_html}
  <div class="hero">
    <div class="hero-left">
      <div class="hero-cat">{cat}</div>
      <div class="hero-title">{titlu}</div>
    </div>
    <div class="hero-right">
      <span class="hl1">{promo_s}</span>
      <span class="hl2">{promo_m}</span>
    </div>
  </div>
  <div class="grid">{cards_html}</div>
  <div class="obs">
    <div class="obs-left">
      <div class="obs-title">OBSERVAȚII</div>
      {obs_html}
    </div>
    <div class="obs-right">
      <div class="obs-brand">{nume}</div>
      Tel: {tel} · {email}<br>
      {web1} · {web2}
    </div>
  </div>
  <div class="foot">
    <span>{footer}</span>
    <span>Pagina {page_num} / {total_pages}</span>
  </div>
</div>'''


# ══════════════════════════════════════════════════════════════
# GENERATE HTML — returnează string-ul HTML
# ══════════════════════════════════════════════════════════════
def generate_html(produse, antet, observatii, style, img_cache,
                  antet_b64="", per_page=6) -> str:
    """Generează HTML-ul flyerului ca string.

    img_cache: dict {filename: data_url_base64}
    antet_b64: data URL al imaginii antet (sau gol pentru fără antet)
    """
    chunk_size = _products_per_page(per_page)

    def chunk(lst, n):
        for i in range(0, len(lst), n):
            yield lst[i:i + n]

    pages = list(chunk(produse, chunk_size))
    total = len(pages)

    pages_html = "\n\n".join(
        build_page(prods, idx + 1, total, antet, observatii, antet_b64, img_cache)
        for idx, prods in enumerate(pages)
    )

    css = build_css(style, per_page=per_page)
    title = antet.get("Titlu produs (mare)", "Flyer")
    nume = antet.get("Nume companie", "")

    return f'''<!DOCTYPE html>
<html lang="ro">
<head>
<meta charset="UTF-8">
<title>Flyer {title} – {nume}</title>
<style>
{css}
</style>
</head>
<body>
{pages_html}

<button class="print-btn" onclick="descarcaPDF()" title="Descarcă ca PDF">
  Descarcă PDF
</button>
<div class="print-hint" id="printHint">
  💡 În fereastra de print, alege <b>»Save as PDF«</b> ca destinație și asigură-te că e bifat <b>»Background graphics«</b>.
</div>

<script>
function descarcaPDF() {{
  var hint = document.getElementById('printHint');
  hint.classList.add('show');
  setTimeout(function() {{
    window.print();
    setTimeout(function() {{ hint.classList.remove('show'); }}, 500);
  }}, 200);
}}
// Ascultă combinația Ctrl+P / Cmd+P pentru hint vizual
window.addEventListener('beforeprint', function() {{
  var hint = document.getElementById('printHint');
  if (hint) hint.classList.remove('show');
}});
</script>
</body>
</html>'''


# ══════════════════════════════════════════════════════════════
# BUILD PPTX — returnează bytes
# ══════════════════════════════════════════════════════════════
def build_pptx(produse, antet, observatii, style, image_files,
               antet_image_bytes: Optional[bytes] = None, per_page=6) -> Optional[bytes]:
    """Generează un fișier .pptx ca bytes. Returnează None dacă python-pptx lipsește.

    image_files: dict {filename: bytes} - fișiere imagine raw
    antet_image_bytes: bytes ale imaginii de antet
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        return None

    # Salvăm temporar imaginile pentru python-pptx (care vrea fișiere de pe disk)
    import tempfile
    tmpdir = tempfile.mkdtemp(prefix="flyer_pptx_")
    image_paths = {}
    for name, data in image_files.items():
        if data:
            p = Path(tmpdir) / name
            p.write_bytes(data)
            image_paths[name] = str(p)

    antet_path = ""
    if antet_image_bytes:
        ap = Path(tmpdir) / "_antet.jpg"
        ap.write_bytes(antet_image_bytes)
        antet_path = str(ap)

    SLIDE_W = Inches(8.27)
    SLIDE_H = Inches(11.69)

    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H
    blank_layout = pres.slide_layouts[6]

    PAGE_PAD   = Inches(0.0)
    HEADER_H   = Inches(1.45)
    HERO_H     = Inches(0.65)
    GRID_TOP   = HEADER_H + HERO_H + Inches(0.05)
    OBS_H      = Inches(0.85)
    FOOT_H     = Inches(0.20)
    GRID_H     = SLIDE_H - GRID_TOP - OBS_H - FOOT_H
    GRID_PAD   = Inches(0.12)
    GRID_GAP   = Inches(0.08)

    PPTX_LAYOUTS = {
        1: (1, 1), 2: (1, 2), 3: (1, 3),
        "3a": (2, 3), 4: (2, 2), 5: (2, 3), 6: (2, 3),
    }
    CARD_COLS, CARD_ROWS = PPTX_LAYOUTS.get(per_page, (2, 3))
    CARD_W = (SLIDE_W - 2*GRID_PAD - (CARD_COLS-1)*GRID_GAP) / CARD_COLS
    CARD_H = (GRID_H - (CARD_ROWS-1)*GRID_GAP) / CARD_ROWS
    PRICE_BAR_H = Inches(0.7)
    IMG_W = CARD_W * 0.46

    BG_PAGE       = hex_to_rgb(style.get("Culoare fundal pagină", DEFAULT_STYLE["Culoare fundal pagină"]))
    BG_RED        = hex_to_rgb(style.get("Culoare principală (roșu)", DEFAULT_STYLE["Culoare principală (roșu)"]))
    BG_GREY       = hex_to_rgb(style.get("Culoare secundară (gri)", DEFAULT_STYLE["Culoare secundară (gri)"]))
    BG_CARD       = hex_to_rgb(style.get("Culoare fundal casetă produs", DEFAULT_STYLE["Culoare fundal casetă produs"]))
    BG_PILL       = hex_to_rgb(style.get("Culoare fundal pill caracteristici", DEFAULT_STYLE["Culoare fundal pill caracteristici"]))
    BORDER_C      = hex_to_rgb(style.get("Culoare bordură", DEFAULT_STYLE["Culoare bordură"]))
    FOOTER_C      = hex_to_rgb(style.get("Culoare footer (jos pagină)", DEFAULT_STYLE["Culoare footer (jos pagină)"]))
    BG_BADGE      = hex_to_rgb(style.get("Culoare fundal badge", DEFAULT_STYLE["Culoare fundal badge"]))
    TXT_BADGE     = hex_to_rgb(style.get("Culoare text badge", DEFAULT_STYLE["Culoare text badge"]))
    NAME_C        = hex_to_rgb(style.get("Culoare denumire produs", DEFAULT_STYLE["Culoare denumire produs"]))
    PILL_LBL_C    = hex_to_rgb(style.get("Culoare caracteristici (etichetă)", DEFAULT_STYLE["Culoare caracteristici (etichetă)"]))
    PILL_VAL_C    = hex_to_rgb(style.get("Culoare caracteristici (valoare)", DEFAULT_STYLE["Culoare caracteristici (valoare)"]))
    COD_C         = hex_to_rgb(style.get("Culoare cod produs", DEFAULT_STYLE["Culoare cod produs"]))
    PRICE_C       = hex_to_rgb(style.get("Culoare preț", DEFAULT_STYLE["Culoare preț"]))
    TVA_C         = hex_to_rgb(style.get("Culoare »Preț cu TVA«", DEFAULT_STYLE["Culoare »Preț cu TVA«"]))
    CAT_C         = hex_to_rgb(style.get("Culoare categorie (sus)", DEFAULT_STYLE["Culoare categorie (sus)"]))
    TITLU_C       = hex_to_rgb(style.get("Culoare titlu produs (mare)", DEFAULT_STYLE["Culoare titlu produs (mare)"]))
    HL1_C         = hex_to_rgb(style.get("Culoare text promoțional mic", DEFAULT_STYLE["Culoare text promoțional mic"]))
    HL2_C         = hex_to_rgb(style.get("Culoare text promoțional mare", DEFAULT_STYLE["Culoare text promoțional mare"]))
    OBS_TITLE_C   = hex_to_rgb(style.get("Culoare titlu observații", DEFAULT_STYLE["Culoare titlu observații"]))
    OBS_TEXT_C    = hex_to_rgb(style.get("Culoare text observații", DEFAULT_STYLE["Culoare text observații"]))
    BRAND_C       = hex_to_rgb(style.get("Culoare nume firmă (footer)", DEFAULT_STYLE["Culoare nume firmă (footer)"]))

    def add_text(slide, x, y, w, h, text, font_name="Arial", size=11,
                 bold=False, color=(0,0,0), align="left", anchor="middle"):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.word_wrap = True
        tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                              "middle": MSO_ANCHOR.MIDDLE}.get(anchor, MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
                       "left": PP_ALIGN.LEFT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = str(text)
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return tb

    def add_rect(slide, x, y, w, h, fill_color=(255,255,255), border=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(*fill_color)
        if border:
            sh.line.color.rgb = RGBColor(*border)
            sh.line.width = Pt(0.75)
        else:
            sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def add_image(slide, path, x, y, w=None, h=None):
        if not path or not Path(path).exists():
            return None
        try:
            kwargs = {}
            if w is not None: kwargs["width"] = w
            if h is not None: kwargs["height"] = h
            return slide.shapes.add_picture(str(path), x, y, **kwargs)
        except Exception:
            return None

    def chunk(lst, n):
        for i in range(0, len(lst), n):
            yield lst[i:i + n]

    chunk_size = _products_per_page(per_page)
    pages = list(chunk(produse, chunk_size))
    total_pages = len(pages)

    for page_idx, page_products in enumerate(pages, start=1):
        slide = pres.slides.add_slide(blank_layout)
        if BG_PAGE != (255, 255, 255):
            add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=BG_PAGE)

        # HEADER
        if antet_path and Path(antet_path).exists():
            add_image(slide, antet_path, 0, 0, w=SLIDE_W, h=HEADER_H)
        add_rect(slide, 0, HEADER_H - Inches(0.025), SLIDE_W, Inches(0.025), fill_color=BG_RED)

        # HERO
        hero_y = HEADER_H
        add_rect(slide, 0, hero_y, SLIDE_W, HERO_H, fill_color=BG_GREY)
        add_rect(slide, SLIDE_W - Inches(1.7), hero_y, Inches(1.7), HERO_H, fill_color=BG_RED)

        add_text(slide, Inches(0.3), hero_y + Inches(0.05), Inches(5.5), Inches(0.22),
                 antet.get("Categorie (sus)", ""),
                 font_name=style.get("Font categorie (sus)", "Arial"),
                 size=int(style.get("Dimensiune categorie (sus)", 9)) + 1,
                 color=CAT_C, anchor="middle")
        add_text(slide, Inches(0.3), hero_y + Inches(0.20), Inches(5.5), HERO_H - Inches(0.25),
                 antet.get("Titlu produs (mare)", ""),
                 font_name=style.get("Font titlu produs (mare)", "Arial"),
                 size=int(style.get("Dimensiune titlu produs (mare)", 28)),
                 bold=True, color=TITLU_C, anchor="top")
        add_text(slide, SLIDE_W - Inches(1.65), hero_y + Inches(0.07), Inches(1.55), Inches(0.20),
                 antet.get("Text promoțional mic", ""),
                 font_name=style.get("Font text promoțional mic", "Arial"),
                 size=int(style.get("Dimensiune text promoțional mic", 9)) + 1,
                 color=HL1_C, align="right", anchor="middle")
        add_text(slide, SLIDE_W - Inches(1.65), hero_y + Inches(0.30), Inches(1.55), Inches(0.32),
                 antet.get("Text promoțional mare", ""),
                 font_name=style.get("Font text promoțional mare", "Arial"),
                 size=int(style.get("Dimensiune text promoțional mare", 17)),
                 bold=True, color=HL2_C, align="right", anchor="middle")

        # GRID PRODUSE
        for idx, p in enumerate(page_products):
            row = idx // CARD_COLS
            col = idx % CARD_COLS
            cx = GRID_PAD + col * (CARD_W + GRID_GAP)
            cy = GRID_TOP + row * (CARD_H + GRID_GAP)
            this_card_w = CARD_W
            this_card_h = CARD_H

            if per_page == 5 and idx == 4:
                cx = GRID_PAD
                this_card_w = SLIDE_W - 2 * GRID_PAD
            elif per_page == "3a":
                if idx == 0:
                    cx = GRID_PAD; cy = GRID_TOP
                    this_card_w = SLIDE_W - 2 * GRID_PAD
                    this_card_h = 2 * CARD_H + GRID_GAP
                else:
                    col_low = idx - 1
                    cx = GRID_PAD + col_low * (CARD_W + GRID_GAP)
                    cy = GRID_TOP + 2 * (CARD_H + GRID_GAP)

            add_rect(slide, cx, cy, this_card_w, this_card_h, fill_color=BG_CARD, border=BORDER_C)

            this_img_w = IMG_W if this_card_w == CARD_W else CARD_W * 0.30
            if per_page == "3a" and idx == 0:
                this_img_w = this_card_w * 0.40

            this_card_top_h = this_card_h - PRICE_BAR_H

            add_rect(slide, cx, cy, this_img_w, this_card_top_h, fill_color=(255,255,255), border=BORDER_C)

            prod_img_path = image_paths.get(p.get("imagine", ""), "")
            if prod_img_path:
                pad = Inches(0.08)
                add_image(slide, prod_img_path, cx + pad, cy + pad, w=this_img_w - 2*pad)

            # Bulina
            bulina_text = (p.get("bulina_text") or "").strip()
            badge_display = ""
            if bulina_text:
                badge_display = bulina_text
            else:
                raw_red = p.get("reducere")
                if raw_red is not None and str(raw_red).strip() != "":
                    try:
                        sv = str(raw_red).strip().replace("%","").replace(",",".")
                        num = float(sv)
                        if 0 < num < 1: num *= 100
                        red_val = int(round(num))
                        if red_val > 0:
                            badge_display = f"-{red_val}%"
                    except Exception:
                        pass

            if badge_display:
                diam = Inches(0.65)
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                    cx + Inches(0.08), cy + Inches(0.08), diam, diam)
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(*BG_RED)
                circle.line.color.rgb = RGBColor(255,255,255)
                circle.line.width = Pt(2)
                tf = circle.text_frame
                tf.margin_left = Pt(0); tf.margin_right = Pt(0)
                tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.word_wrap = True
                pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
                run = pp.add_run()
                run.text = badge_display
                run.font.name = style.get("Font preț", "Arial")
                base_size = int(style.get("Dimensiune cifră reducere", 24))
                if "%" in badge_display:
                    run.font.size = Pt(max(9, base_size - 4))
                else:
                    text_len = len(badge_display)
                    if text_len <= 3:    run.font.size = Pt(max(10, base_size - 6))
                    elif text_len <= 5:  run.font.size = Pt(max(8, base_size - 10))
                    else:                run.font.size = Pt(max(7, base_size - 14))
                run.font.bold = True
                run.font.color.rgb = RGBColor(255,255,255)

            # Info
            info_x = cx + this_img_w + Inches(0.10)
            info_w = this_card_w - this_img_w - Inches(0.15)

            add_text(slide, info_x, cy + Inches(0.10), info_w, Inches(0.40),
                     (p.get("denumire","") or "").upper(),
                     font_name=style.get("Font denumire produs", "Arial"),
                     size=int(style.get("Dimensiune denumire produs", 18)) - 4,
                     bold=True, color=NAME_C, anchor="top")

            add_text(slide, info_x, cy + Inches(0.52), info_w, Inches(0.18),
                     f"COD: {p.get('cod','')}",
                     font_name=style.get("Font cod produs", "Courier New"),
                     size=int(style.get("Dimensiune cod produs", 13)) - 3,
                     color=COD_C, anchor="top")

            pill_y = cy + Inches(0.72)
            pill_h = Inches(0.18)
            for a in p.get("attrs", [])[:3]:
                if not a.get("label") or not a.get("value"):
                    continue
                add_rect(slide, info_x, pill_y, info_w - Inches(0.05), pill_h,
                         fill_color=BG_PILL, border=BORDER_C)
                tb = slide.shapes.add_textbox(info_x + Inches(0.05), pill_y, info_w - Inches(0.10), pill_h)
                tf = tb.text_frame
                tf.margin_left = Pt(0); tf.margin_right = Pt(0)
                tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
                r1 = pp.add_run(); r1.text = f"{a['label']}: "
                r1.font.name = style.get("Font caracteristici", "Arial")
                r1.font.size = Pt(int(style.get("Dimensiune caracteristici", 13)) - 4)
                r1.font.color.rgb = RGBColor(*PILL_LBL_C)
                r2 = pp.add_run(); r2.text = str(a["value"])
                r2.font.name = style.get("Font caracteristici", "Arial")
                r2.font.size = Pt(int(style.get("Dimensiune caracteristici", 13)) - 4)
                r2.font.bold = True
                r2.font.color.rgb = RGBColor(*PILL_VAL_C)
                pill_y += pill_h + Inches(0.04)

            # Banda preț
            price_y = cy + this_card_top_h
            add_rect(slide, cx, price_y, this_card_w, PRICE_BAR_H, fill_color=BG_RED)

            pret_int = format_price(p.get("pret"))
            pret_redus = format_price(p.get("pret_redus"))
            has_red = p.get("pret_redus") is not None and str(p.get("pret_redus","")).strip() != "" and pret_redus

            moneda = (p.get("moneda") or "").strip()
            um = (p.get("um") or "").strip()
            sufix_complet = ""
            if moneda: sufix_complet += f" {moneda}"
            if um:     sufix_complet += f" / {um}"
            sufix_doar_moneda = f" {moneda}" if moneda else ""

            price_scale = 1.0
            if per_page == "3a":
                price_scale = 1.4 if idx == 0 else 0.7
            base_price_size = int(style.get("Dimensiune preț", 32)) - 6

            price_x = cx + Inches(0.15)
            price_w = this_card_w * 0.55

            if has_red:
                tb = slide.shapes.add_textbox(price_x, price_y + Inches(0.05), price_w, Inches(0.22))
                tf = tb.text_frame
                tf.margin_left = Pt(0); tf.margin_right = Pt(0)
                tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
                run = pp.add_run(); run.text = f"{pret_int}{sufix_doar_moneda}"
                run.font.name = style.get("Font preț", "Arial")
                run.font.size = Pt(max(8, int((int(style.get("Dimensiune preț vechi (tăiat)", 16)) - 3) * price_scale)))
                run.font.color.rgb = RGBColor(*hex_to_rgb(style.get("Culoare preț vechi (tăiat)", "FFB0B0")))
                rPr = run._r.get_or_add_rPr()
                rPr.set("strike", "sngStrike")
                add_text(slide, price_x, price_y + Inches(0.30), price_w, Inches(0.40),
                         f"{pret_redus}{sufix_complet}",
                         font_name=style.get("Font preț", "Arial"),
                         size=max(10, int(base_price_size * price_scale)),
                         bold=True, color=PRICE_C, anchor="top")
            else:
                add_text(slide, price_x, price_y, price_w, PRICE_BAR_H,
                         f"{pret_int}{sufix_complet}",
                         font_name=style.get("Font preț", "Arial"),
                         size=max(10, int(base_price_size * price_scale)),
                         bold=True, color=PRICE_C, anchor="middle")

            tva_text = (p.get("tva_text") or "").strip()
            if tva_text:
                add_text(slide, cx + this_card_w * 0.58, price_y, this_card_w * 0.40, PRICE_BAR_H,
                         tva_text.upper(),
                         font_name=style.get("Font preț", "Arial"),
                         size=int(style.get("Dimensiune »Preț cu TVA«", 13)) - 4,
                         color=TVA_C, align="right", anchor="middle")

        # OBSERVAȚII
        obs_y = SLIDE_H - OBS_H - FOOT_H
        add_rect(slide, 0, obs_y, SLIDE_W, OBS_H, fill_color=BG_PILL)
        add_text(slide, Inches(0.3), obs_y + Inches(0.08), Inches(2), Inches(0.25),
                 "OBSERVAȚII",
                 font_name=style.get("Font observații", "Arial"),
                 size=int(style.get("Dimensiune observații", 12)) + 1,
                 bold=True, color=OBS_TITLE_C, anchor="top")
        for i, obs_line in enumerate(observatii):
            add_text(slide, Inches(0.3), obs_y + Inches(0.30) + i * Inches(0.18),
                     Inches(4.5), Inches(0.18), f"▪ {obs_line}",
                     font_name=style.get("Font observații", "Arial"),
                     size=int(style.get("Dimensiune observații", 12)) - 1,
                     color=OBS_TEXT_C, anchor="top")
        cx_right = SLIDE_W - Inches(3.0)
        add_text(slide, cx_right, obs_y + Inches(0.08), Inches(2.7), Inches(0.25),
                 antet.get("Nume companie", ""),
                 font_name=style.get("Font observații", "Arial"),
                 size=int(style.get("Dimensiune observații", 12)) + 2,
                 bold=True, color=BRAND_C, align="right", anchor="top")
        add_text(slide, cx_right, obs_y + Inches(0.32), Inches(2.7), Inches(0.18),
                 f"Tel: {antet.get('Telefon','')} · {antet.get('Email','')}",
                 font_name=style.get("Font observații", "Arial"),
                 size=int(style.get("Dimensiune observații", 12)) - 2,
                 color=OBS_TEXT_C, align="right", anchor="top")
        add_text(slide, cx_right, obs_y + Inches(0.50), Inches(2.7), Inches(0.18),
                 f"{antet.get('Website 1','')} · {antet.get('Website 2','')}",
                 font_name=style.get("Font observații", "Arial"),
                 size=int(style.get("Dimensiune observații", 12)) - 2,
                 color=OBS_TEXT_C, align="right", anchor="top")

        # FOOTER
        foot_y = SLIDE_H - FOOT_H
        add_rect(slide, 0, foot_y, SLIDE_W, FOOT_H, fill_color=FOOTER_C)
        add_text(slide, Inches(0.3), foot_y, Inches(5), FOOT_H,
                 antet.get("Footer legal", ""),
                 font_name="Arial", size=6, color=(180,180,180), anchor="middle")
        add_text(slide, SLIDE_W - Inches(1.5), foot_y, Inches(1.2), FOOT_H,
                 f"Pagina {page_idx} / {total_pages}",
                 font_name="Arial", size=6, color=(180,180,180), align="right", anchor="middle")

    # Salvăm într-un buffer și returnăm bytes
    out_buf = io.BytesIO()
    pres.save(out_buf)
    out_buf.seek(0)

    # Curățăm folderul temp
    import shutil
    shutil.rmtree(tmpdir, ignore_errors=True)

    return out_buf.getvalue()


# ══════════════════════════════════════════════════════════════
# HTML → PDF (folosind WeasyPrint)
# ══════════════════════════════════════════════════════════════
def html_to_pdf(html_content: str) -> Optional[bytes]:
    """Convertește un string HTML într-un fișier PDF (bytes).

    Folosește WeasyPrint - bibliotecă Python pură care suportă CSS modern
    (grid, flexbox, Google Fonts, @page rules).

    Returnează None dacă WeasyPrint nu este instalat sau apare o eroare.
    """
    try:
        from weasyprint import HTML
    except ImportError:
        return None

    try:
        # Generăm PDF în memorie (fără fișier temporar)
        pdf_bytes = HTML(string=html_content).write_pdf()
        return pdf_bytes
    except Exception as e:
        # Loggam eroarea dar nu o aruncăm - returnăm None ca să poată app-ul
        # să-și dea seama că PDF-ul nu e disponibil
        print(f"WeasyPrint error: {e}")
        return None
